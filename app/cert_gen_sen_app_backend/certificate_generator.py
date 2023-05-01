import glob
import os
from .models import *
from django.http import JsonResponse
from django.core.mail import EmailMessage
from rest_framework.views import APIView
from django.conf import settings
from collections import OrderedDict
from python_pptx_text_replacer import TextReplacer
import tqdm
import qrcode
from pptx import Presentation
from PIL import Image
import smtplib
import ssl
from twilio.rest import Client
from decouple import config
from email import encoders
from email.mime.base import MIMEBase
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from pypdf import PdfMerger


# PDF merger
def pdf_certificate_merger(pdf_list, event_slug):
    merger = PdfMerger()

    for pdf in pdf_list:
        merger.append(pdf)

    merger.write(
        f'./cert_gen_sen_app_backend/certificate_data/merged-certificates/Certificates - {event_slug}.pdf')
    merger.close()

    Event.objects.update(
        certificates_file=f'./cert_gen_sen_app_backend/certificate_data/merged-certificates/Certificates - {event_slug}.pdf', certificate_file_name=f'Certificates - {event_slug}.pdf')


# Sending message to participant
def send_message(participant_name, phone, senders_phone):
    account_sid = config("TWILIO_SID")
    auth_token = config("TWILIO_AUTH_TOKEN")
    client = Client(account_sid, auth_token)
    client.messages.create(
        body=f"Dear {participant_name}, thankyou for participating in the event/contest. Your certificate will be delivered to you via e-mail. Check your email.",
        from_="+15855951968",
        # from_=senders_phone,
        to=phone
    )
    return "SENT"


# Sending mail to participant
def send_mail(subject, body, email_to, certificate_file, senders_email, senders_password):
    email_from = senders_email
    password = senders_password
    try:
        message = MIMEMultipart()
        message['From'] = email_from
        message['To'] = email_to
        message['Subject'] = subject
        message['Bcc'] = email_to

        message.attach(MIMEText(body, 'plain'))
        file = certificate_file

        with open(file, 'rb') as attachment:
            part = MIMEBase('application', 'octet-stream')
            part.set_payload(attachment.read())

        encoders.encode_base64(part)

        part.add_header('Content-Disposition',
                        f"attachment;filename={file.replace('./cert_gen_sen_app_backend/certificate_data/participants-certificates/', '')}")

        message.attach(part)
        text = message.as_string()

        context = ssl.create_default_context()
        with smtplib.SMTP_SSL('smtp.gmail.com', 465, context=context) as server:
            server.login(email_from, password)
            server.sendmail(email_from, email_to, text)

        return "SENT"
    except Exception as e:
        print(e)


# Sending mail to each participant
def sendMail(subject, message, email_to, certificate_file):
    try:
        email_form = settings.EMAIL_HOST_USER
        certificate = EmailMessage(subject, message, email_form, [email_to])
        certificate.attach_file(certificate_file)
        certificate.send()
        return "SENT"
    except Exception as e:
        print(e)


# Certificate directory cleaner
def cleanUp():
    participant_files = "../app/cert_gen_sen_app_backend/certificate_data/participants-certificates"
    participant_filelist = glob.glob(os.path.join(participant_files, "*"))
    for f in participant_filelist:
        os.remove(f)

    merit_files = "../app/cert_gen_sen_app_backend/certificate_data/merit-certificates"
    merit_filelist = glob.glob(os.path.join(merit_files, "*"))
    for f in merit_filelist:
        os.remove(f)


# Placing QR code in PPT
def place_qrcode(pptx_path, qrcode_path, replace_str):
    pptx_file = Presentation(pptx_path)

    for slide in pptx_file.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text.find(replace_str) != -1:
                    slide.shapes.add_picture(
                        qrcode_path, shape.left, shape.top)
    pptx_file.save(pptx_path)


# Set size of QR code
def resize_qrcode(qrcode_path):
    img = Image.open(qrcode_path)
    width = 100
    height = 100
    resize = img.resize((width, height), Image.NEAREST)
    resize.save(qrcode_path)


# Generate QR code
def generate_qrcode(stu_name, stu_id, cert_id, eve_name, eve_department, eve_date):
    qr_code = qrcode.QRCode(
        version=5, box_size=2, border=1)
    qr_code.add_data(
        f"Participant Name - {stu_name}\nParticipant Id - {stu_id}\nEvent Name - {eve_name}\nEvent Department - {eve_department}\nEvent Date - {eve_date}\nCertificate Id - {cert_id}")
    type(qr_code)
    qr_code.make(fit=True)
    img = qr_code.make_image(fill_color="black", back_color="white")
    img.save(
        f'./cert_gen_sen_app_backend/certificate_data/qr-code/{cert_id} - {stu_name}.png')

    # resize_qrcode(
    #     f'./cert_gen_sen_app_backend/certificate_data/qr-code/{cert_id} - {stu_name}.png')

    return f'./cert_gen_sen_app_backend/certificate_data/qr-code/{cert_id} - {stu_name}.png'


# Generate certificates for all participants type
def generate_participant_certificate(stu_name, cert_id, qrcode_path, completion_certificate_path):
    replacer = TextReplacer(f'../app{completion_certificate_path}',
                            slides="", tables=True, charts=True, textframes=True)

    replacer.replace_text(
        [("{{StudentName}}", stu_name), ("{{UID}}", cert_id)])

    replacer.write_presentation_to_file(
        f'./cert_gen_sen_app_backend/certificate_data/ppt-certificates/{cert_id} - {stu_name}.pptx')

    pptx_path = f'./cert_gen_sen_app_backend/certificate_data/ppt-certificates/{cert_id} - {stu_name}.pptx'

    place_qrcode(pptx_path, qrcode_path, "{{QR}}")

    path = './cert_gen_sen_app_backend/certificate_data/ppt-certificates'
    ext = 'pptx'

    files = [f for f in glob.glob(
        path + "/**/*.{}".format(ext), recursive=True)]

    for f in tqdm.tqdm(files):
        command = "unoconv -f pdf \"{}\"".format(f)
        move_file = "mv ./cert_gen_sen_app_backend/certificate_data/ppt-certificates/*.pdf ./cert_gen_sen_app_backend/certificate_data/participants-certificates/"
        os.system(command)
        os.system(move_file)

    return f'./cert_gen_sen_app_backend/certificate_data/participants-certificates/{cert_id} - {stu_name}.pdf'


# Generate certificates for merit participants
def generate_merit_certificate(stu_name, cert_id, rank, qrcode_path, merit_certificate_path):
    replacer = TextReplacer(f'../app{merit_certificate_path}',
                            slides="", tables=True, charts=True, textframes=True)

    if (rank == "1"):
        replacer.replace_text(
            [("{{StudentName}}", stu_name), ("{{UID}}", cert_id), ("{{POS}}", f"{rank} st")])
    elif (rank == "2"):
        replacer.replace_text(
            [("{{StudentName}}", stu_name), ("{{UID}}", cert_id), ("{{POS}}", f"{rank} nd")])
    elif (rank == "3"):
        replacer.replace_text(
            [("{{StudentName}}", stu_name), ("{{UID}}", cert_id), ("{{POS}}", f"{rank} rd")])

    replacer.write_presentation_to_file(
        f'./cert_gen_sen_app_backend/certificate_data/ppt-certificates/{cert_id} - {stu_name}.pptx')

    pptx_path = f'./cert_gen_sen_app_backend/certificate_data/ppt-certificates/{cert_id} - {stu_name}.pptx'

    place_qrcode(pptx_path, qrcode_path, "{{QR}}")

    path = './cert_gen_sen_app_backend/certificate_data/ppt-certificates'
    ext = 'pptx'

    files = [f for f in glob.glob(
        path + "/**/*.{}".format(ext), recursive=True)]

    for f in tqdm.tqdm(files):
        command = "unoconv -f pdf \"{}\"".format(f)
        move_file = "mv ./cert_gen_sen_app_backend/certificate_data/ppt-certificates/*.pdf ./cert_gen_sen_app_backend/certificate_data/merit-certificates/"
        os.system(command)
        os.system(move_file)

    return f'./cert_gen_sen_app_backend/certificate_data/merit-certificates/{cert_id} - {stu_name}.pdf'


# Generate certificates for all
class GenerateCertificate(APIView):
    def post(self, request, slug):
        try:
            obj_event = Event.objects.get(slug=slug)
            event = Event.objects.filter(slug=slug)
            participants = Participant.objects.filter(event=obj_event)

            completion_certificate_path = request.data['completion']
            merit_certificate_path = request.data['merit']

            senders_email = SendersCredentials.objects.filter(
                user=request.user.id).values('senders_email').first()['senders_email']
            senders_password = SendersCredentials.objects.filter(
                user=request.user.id).values('senders_password').first()['senders_password']
            senders_phone = SendersCredentials.objects.filter(
                user=request.user.id).values('senders_phone').first()['senders_phone']

            stu_data = []
            eve_data = OrderedDict()

            merge_certificate_path = []

            for stu in participants:
                if stu.certificate_sent_status == False:
                    if stu.certificate_status == 'T' or stu.certificate_status == '1' or stu.certificate_status == '2' or stu.certificate_status == '3':
                        stu_data.append(dict(id=stu.id, participant_name=stu.participant_name, participant_id=stu.participant_id, email=stu.email, phone=stu.phone,
                                             certificate_status=stu.certificate_status, certificate_id=stu.certificate_id))

            for eve in event:
                eve_data['event_name'] = eve.event_name
                eve_data['event_department'] = eve.event_department
                eve_data['from_date'] = eve.from_date.strftime('%d-%m-%Y')
                eve_data['slug'] = eve.slug

            for stu in stu_data:
                if stu['certificate_status'] == "1" or stu['certificate_status'] == "2" or stu['certificate_status'] == "3":
                    qrcode_path = generate_qrcode(
                        stu["participant_name"], stu["participant_id"], stu["certificate_id"], eve_data["event_name"], eve_data["event_department"], eve_data["from_date"])
                    certificate_path = generate_merit_certificate(
                        stu['participant_name'], stu['certificate_id'], stu['certificate_status'], qrcode_path, merit_certificate_path)
                    send_message(stu["participant_name"],
                                 stu["phone"], senders_phone)
                    send_certificate = send_mail("Certificate of Participation",
                                                 "Thank you for participanting in the Event/Contest", stu["email"], certificate_path, senders_email, senders_password)

                    if send_certificate == "SENT":
                        Participant.objects.filter(
                            id=stu['id']).update(certificate_sent_status=True)

                    merge_certificate_path.append(certificate_path)
                else:
                    qrcode_path = generate_qrcode(
                        stu["participant_name"], stu["participant_id"], stu["certificate_id"], eve_data["event_name"], eve_data["event_department"], eve_data["from_date"])
                    certificate_path = generate_participant_certificate(
                        stu["participant_name"], stu["certificate_id"], qrcode_path, completion_certificate_path)
                    send_message(stu["participant_name"],
                                 stu["phone"], senders_phone)
                    send_certificate = send_mail("Certificate of Participation",
                                                 "Thank you for participanting in the Event/Contest", stu["email"], certificate_path, senders_email, senders_password)

                    if send_certificate == "SENT":
                        Participant.objects.filter(
                            id=stu['id']).update(certificate_sent_status=True)

                    merge_certificate_path.append(certificate_path)

            pdf_certificate_merger(merge_certificate_path, eve_data["slug"])

            return JsonResponse("Certificate generated and sended successfully", safe=False)
        except:
            return JsonResponse("Some problem occured while sending", safe=False)


# Generate certificates for specific participant
class GenerateCertificateById(APIView):
    def post(self, request, slug, pk):
        participant = Participant.objects.filter(id=pk)
        event = Event.objects.filter(slug=slug)

        stu_data = OrderedDict()
        eve_data = OrderedDict()

        senders_email = SendersCredentials.objects.filter(
            user=request.user.id).values('senders_email').first()['senders_email']
        senders_password = SendersCredentials.objects.filter(
            user=request.user.id).values('senders_password').first()['senders_password']
        senders_phone = SendersCredentials.objects.filter(
            user=request.user.id).values('senders_phone').first()['senders_phone']

        completion_certificate_path = request.data['completion']
        merit_certificate_path = request.data['merit']

        for stu in participant:
            stu_data['participant_name'] = stu.participant_name
            stu_data['participant_id'] = stu.participant_id
            stu_data['email'] = stu.email
            stu_data['phone'] = stu.phone
            stu_data['certificate_status'] = stu.certificate_status
            stu_data['certificate_id'] = stu.certificate_id

        for eve in event:
            eve_data['event_name'] = eve.event_name
            eve_data['event_department'] = eve.event_department
            eve_data['from_date'] = eve.from_date.strftime('%d-%m-%Y')

        if stu_data["certificate_status"] == 'F' or stu_data["certificate_status"] == None:
            return JsonResponse("This participant is not eligible for certificate", safe=False)
        elif stu_data["certificate_status"] == '1' or stu_data["certificate_status"] == '2' or stu_data["certificate_status"] == '3':
            certificate_path = generate_merit_certificate(
                stu_data["name"], stu_data["certificate_id"], stu_data["certificate_status"], qrcode_path, merit_certificate_path)
        else:
            qrcode_path = generate_qrcode(
                stu_data["participant_name"], stu_data["participant_id"], stu_data["certificate_id"], eve_data["event_name"], eve_data["event_department"], eve_data["from_date"])
            certificate_path = generate_participant_certificate(
                stu_data["participant_name"], stu_data["certificate_id"], qrcode_path, completion_certificate_path)
            send_message(stu_data["participant_name"],
                         stu_data["phone"], senders_phone)
            send_certificate = send_mail("Certificate of Participation",
                                         "Thank you for participanting in the Event/Contest", stu_data["email"], certificate_path, senders_email, senders_password)

            if send_certificate == "SENT":
                Participant.objects.filter(
                    id=pk).update(certificate_sent_status=True)
                return JsonResponse("Certificate generated and sended successfully", safe=False)

        return JsonResponse("Some problem occured while sending", safe=False)
